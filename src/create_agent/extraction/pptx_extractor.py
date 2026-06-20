"""Extract text from .pptx files using python-pptx."""

from __future__ import annotations

from pathlib import Path


def extract_pptx_text(file_path: str | Path) -> str:
    """Extract all text from a .pptx file.

    Reads text from all slides, including shapes, text frames, and tables.

    Args:
        file_path: Path to the .pptx file.

    Returns:
        Extracted text content with slide separators.

    Raises:
        ImportError: If python-pptx is not installed.
        ValueError: If the file is not a valid .pptx file.
    """
    from pptx import Presentation

    path = Path(file_path)
    if not path.suffix.lower() == ".pptx":
        raise ValueError(f"Not a .pptx file: {path}")

    prs = Presentation(str(path))

    parts: list[str] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_texts: list[str] = []

        for shape in slide.shapes:
            # Text frames (titles, text boxes)
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)

            # Tables
            if shape.has_table:
                for row in shape.table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells]
                    slide_texts.append(" | ".join(row_texts))

        if slide_texts:
            parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))

    return "\n\n".join(parts)
